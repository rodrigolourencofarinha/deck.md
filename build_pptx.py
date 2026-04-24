#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
import json
import re
import sys

import yaml


BLUE = RGBColor(21, 56, 101)
TEXT = RGBColor(33, 33, 33)
MUTED = RGBColor(68, 68, 68)
PANEL = RGBColor(235, 242, 250)
PANEL_BORDER = RGBColor(210, 222, 238)
WHITE = RGBColor(255, 255, 255)
SKILL_DIR = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = SKILL_DIR / "assets" / "RLF_PPT_Template_v1.pptx"


def add_takeaway_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    bullets: list[dict[str, str]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    panel = slide.shapes.add_shape(1, Inches(0.45), Inches(0.35), Inches(12.35), Inches(6.45))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    top_y = 1.18
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.12), Inches(11.2), Inches(0.68))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(15)
        r.font.italic = True
        r.font.color.rgb = MUTED
        top_y = 1.82

    bullets = bullets or []
    if bullets:
        bullets_box = slide.shapes.add_textbox(Inches(0.9), Inches(top_y), Inches(11.0), Inches(4.05 if closing else 4.7))
        tf = bullets_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        first = True
        for item in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            p.bullet = True
            p.space_after = Pt(6)
            p.line_spacing = 1.05

            head = item.get('lead', '').strip()
            body = item.get('body', '').strip()

            if head:
                r1 = p.add_run()
                r1.text = head
                r1.font.name = 'Aptos'
                r1.font.size = Pt(15)
                r1.font.bold = True
                r1.font.color.rgb = TEXT
            if body:
                r2 = p.add_run()
                prefix = ' ' if head and not body.startswith(' ') else ''
                r2.text = f'{prefix}{body}'
                r2.font.name = 'Aptos'
                r2.font.size = Pt(14)
                r2.font.color.rgb = TEXT

    if closing:
        shape = slide.shapes.add_shape(1, Inches(0.8), Inches(6.08), Inches(11.1), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shape.line.color.rgb = PANEL_BORDER
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.bold = False
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_agenda_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    sections: list[dict[str, str]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    panel = slide.shapes.add_shape(1, Inches(0.45), Inches(0.35), Inches(12.35), Inches(6.45))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    top_y = 1.18
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.12), Inches(11.2), Inches(0.55))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.color.rgb = MUTED
        top_y = 1.75

    sections = sections or []
    start_y = top_y
    box_h = 0.84
    gap = 0.16
    accent_x = 0.95
    content_x = 1.35

    for idx, item in enumerate(sections[:5], start=1):
        y = start_y + (idx - 1) * (box_h + gap)
        accent = slide.shapes.add_shape(1, Inches(accent_x), Inches(y + 0.02), Inches(0.26), Inches(box_h - 0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = BLUE
        accent.line.color.rgb = BLUE

        card = slide.shapes.add_shape(1, Inches(content_x), Inches(y), Inches(10.35), Inches(box_h))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = PANEL_BORDER

        num_box = slide.shapes.add_textbox(Inches(content_x + 0.18), Inches(y + 0.12), Inches(0.55), Inches(0.36))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f'{idx:02d}'
        r.font.name = 'Aptos'
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = BLUE

        head_box = slide.shapes.add_textbox(Inches(content_x + 0.8), Inches(y + 0.09), Inches(8.9), Inches(0.28))
        tf = head_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item.get('title', '')
        r.font.name = 'Aptos'
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT

        body = item.get('body', '').strip()
        if body:
            body_box = slide.shapes.add_textbox(Inches(content_x + 0.8), Inches(y + 0.4), Inches(8.9), Inches(0.22))
            tf = body_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = body
            r.font.name = 'Aptos'
            r.font.size = Pt(11)
            r.font.color.rgb = MUTED

    if closing:
        shape = slide.shapes.add_shape(1, Inches(0.8), Inches(6.12), Inches(11.1), Inches(0.42))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shape.line.color.rgb = PANEL_BORDER
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_process_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    phases: list[dict[str, str]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    subtitle_y = 1.02
    track_y = 1.9
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.08), Inches(11.2), Inches(0.42))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED
        track_y = 2.02

    phases = phases or []
    phases = phases[:6]
    count = max(len(phases), 1)
    left = 0.95
    right = 12.0
    track_w = right - left
    step_gap = 0.12
    box_w = (track_w - step_gap * (count - 1)) / count

    spine = slide.shapes.add_shape(1, Inches(left), Inches(track_y + 0.72), Inches(track_w), Inches(0.1))
    spine.fill.solid()
    spine.fill.fore_color.rgb = PANEL_BORDER
    spine.line.color.rgb = PANEL_BORDER

    for idx, phase in enumerate(phases, start=1):
        x = left + (idx - 1) * (box_w + step_gap)
        circle = slide.shapes.add_shape(1, Inches(x + box_w / 2 - 0.18), Inches(track_y + 0.56), Inches(0.36), Inches(0.36))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.color.rgb = BLUE

        ctf = circle.text_frame
        ctf.clear()
        p = ctf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(idx)
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE

        card = slide.shapes.add_shape(1, Inches(x), Inches(track_y + 1.05), Inches(box_w), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = PANEL_BORDER

        head_box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(track_y + 1.18), Inches(box_w - 0.24), Inches(0.45))
        tf = head_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = phase.get('title', '')
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = TEXT

        body = phase.get('body', '').strip()
        if body:
            body_box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(track_y + 1.67), Inches(box_w - 0.24), Inches(1.1))
            tf = body_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = body
            r.font.name = 'Aptos'
            r.font.size = Pt(11)
            r.font.color.rgb = MUTED

    if closing:
        footer = slide.shapes.add_shape(1, Inches(0.9), Inches(5.85), Inches(11.2), Inches(0.48))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(248, 250, 252)
        footer.line.color.rgb = PANEL_BORDER
        tf = footer.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_matrix_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    x_axis: str | None = None,
    y_axis: str | None = None,
    quadrants: list[dict[str, str]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.2), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    matrix_top = 1.35
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.02), Inches(11.2), Inches(0.38))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED
        matrix_top = 1.65

    matrix_left = 1.45
    matrix_w = 9.9
    matrix_h = 4.2

    outer = slide.shapes.add_shape(1, Inches(matrix_left), Inches(matrix_top), Inches(matrix_w), Inches(matrix_h))
    outer.fill.background()
    outer.line.color.rgb = PANEL_BORDER

    vline = slide.shapes.add_shape(1, Inches(matrix_left + matrix_w / 2), Inches(matrix_top), Inches(0.05), Inches(matrix_h))
    vline.fill.solid()
    vline.fill.fore_color.rgb = PANEL_BORDER
    vline.line.color.rgb = PANEL_BORDER

    hline = slide.shapes.add_shape(1, Inches(matrix_left), Inches(matrix_top + matrix_h / 2), Inches(matrix_w), Inches(0.05))
    hline.fill.solid()
    hline.fill.fore_color.rgb = PANEL_BORDER
    hline.line.color.rgb = PANEL_BORDER

    colors = [RGBColor(248, 250, 252), RGBColor(240, 245, 252), RGBColor(245, 249, 253), RGBColor(237, 243, 250)]
    positions = [
        (matrix_left, matrix_top),
        (matrix_left + matrix_w / 2 + 0.025, matrix_top),
        (matrix_left, matrix_top + matrix_h / 2 + 0.025),
        (matrix_left + matrix_w / 2 + 0.025, matrix_top + matrix_h / 2 + 0.025),
    ]
    quad_w = matrix_w / 2 - 0.025
    quad_h = matrix_h / 2 - 0.025

    quadrants = (quadrants or [])[:4]
    for idx in range(4):
        x, y = positions[idx]
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(quad_w), Inches(quad_h))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[idx]
        box.line.color.rgb = WHITE

        item = quadrants[idx] if idx < len(quadrants) else {}
        head_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(quad_w - 0.3), Inches(0.38))
        tf = head_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item.get('title', '')
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = BLUE

        body = item.get('body', '').strip()
        if body:
            body_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.62), Inches(quad_w - 0.3), Inches(1.25))
            tf = body_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = body
            r.font.name = 'Aptos'
            r.font.size = Pt(11)
            r.font.color.rgb = TEXT

    if x_axis:
        x_box = slide.shapes.add_textbox(Inches(matrix_left + 3.1), Inches(matrix_top + matrix_h + 0.12), Inches(3.9), Inches(0.3))
        tf = x_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = x_axis
        r.font.name = 'Aptos'
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED

    if y_axis:
        y_box = slide.shapes.add_textbox(Inches(0.55), Inches(matrix_top + 1.35), Inches(0.7), Inches(1.7))
        tf = y_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = y_axis
        r.font.name = 'Aptos'
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED

    if closing:
        footer = slide.shapes.add_shape(1, Inches(0.9), Inches(6.1), Inches(11.2), Inches(0.42))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(248, 250, 252)
        footer.line.color.rgb = PANEL_BORDER
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_bar_column_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    bars: list[dict[str, str | int | float]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    chart_top = 1.4
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.08), Inches(11.2), Inches(0.38))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED
        chart_top = 1.72

    bars = bars or []
    max_val = max([float(b.get('value', 0)) for b in bars], default=1.0) or 1.0
    base_x = 1.35
    label_w = 1.9
    bar_x = base_x + label_w
    max_bar_w = 8.2
    row_h = 0.72

    for idx, bar in enumerate(bars[:6]):
        y = chart_top + idx * row_h
        label_box = slide.shapes.add_textbox(Inches(base_x), Inches(y + 0.14), Inches(label_w - 0.15), Inches(0.26))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(bar.get('label', ''))
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = TEXT

        value = float(bar.get('value', 0))
        width = max(0.65, max_bar_w * value / max_val)
        shape = slide.shapes.add_shape(1, Inches(bar_x), Inches(y + 0.08), Inches(width), Inches(0.34))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.color.rgb = BLUE

        value_box = slide.shapes.add_textbox(Inches(bar_x + width + 0.1), Inches(y + 0.11), Inches(0.8), Inches(0.24))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f'{value:g}'
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = BLUE

    if closing:
        footer = slide.shapes.add_shape(1, Inches(0.9), Inches(6.1), Inches(11.2), Inches(0.42))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(248, 250, 252)
        footer.line.color.rgb = PANEL_BORDER
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_section_divider_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    section_label: str | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    bg = slide.shapes.add_shape(1, Inches(0.45), Inches(0.35), Inches(12.35), Inches(6.45))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PANEL
    bg.line.color.rgb = PANEL

    if section_label:
        label_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.2), Inches(2.2), Inches(0.32))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = section_label.upper()
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = BLUE

    title_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.72), Inches(10.6), Inches(1.15))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = BLUE

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.98), Inches(3.08), Inches(10.3), Inches(0.45))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT

    if closing:
        box = slide.shapes.add_textbox(Inches(0.98), Inches(5.55), Inches(10.5), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.italic = True
        r.font.color.rgb = MUTED


def add_quote_slide(
    prs: Presentation,
    *,
    title: str,
    quote: str,
    subtitle: str | None = None,
    attribution: str | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = BLUE

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.2), Inches(0.45))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED

    quote_card = slide.shapes.add_shape(1, Inches(1.1), Inches(1.75), Inches(10.9), Inches(3.05))
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = PANEL
    quote_card.line.color.rgb = PANEL_BORDER

    q_box = slide.shapes.add_textbox(Inches(1.45), Inches(2.15), Inches(10.2), Inches(2.0))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f'“{quote}”'
    r.font.name = 'Aptos'
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = BLUE

    if attribution:
        a_box = slide.shapes.add_textbox(Inches(7.8), Inches(4.3), Inches(3.4), Inches(0.35))
        tf = a_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = attribution
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = MUTED

    if closing:
        box = slide.shapes.add_textbox(Inches(1.1), Inches(5.3), Inches(10.9), Inches(0.55))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT


def load_structured_file(path: Path):
    text = path.read_text(encoding='utf-8')
    if path.suffix.lower() in {'.yaml', '.yml', '.md'}:
        data = yaml.safe_load(text)
        if not isinstance(data, dict):
            raise SystemExit(f'spec root must be a mapping: {path}')
        return data
    return json.loads(text)


def slugify(value: str) -> str:
    value = value.lower().strip()
    value = re.sub(r'[^a-z0-9]+', '-', value)
    return value.strip('-') or 'deck'


def normalize_mode(deck_mode: str | None, slide: dict) -> str:
    deck_mode = (deck_mode or '').strip().lower()
    slide_mode = str(slide.get('slide_mode') or '').strip().lower()
    if deck_mode == 'mixed':
        if slide_mode not in {'ppt-shapes', 'designer-mode'}:
            raise SystemExit(f"mixed decks require slide_mode per slide; missing on slide id={slide.get('id')}")
        return slide_mode
    if deck_mode in {'ppt-shapes', 'designer-mode'}:
        return deck_mode
    return slide_mode or 'ppt-shapes'


def storyline_takeaway(spec: dict, slide_id: int | str | None) -> str | None:
    for item in spec.get('Storyline', []) or []:
        if not isinstance(item, dict):
            continue
        if str(item.get('slide')) == str(slide_id):
            return item.get('key_takeaway')
    return None


def infer_archetype(slide: dict) -> str:
    template = str(slide.get('template') or '').strip().lower()
    layout = str(slide.get('layout') or '').strip().lower()
    candidates = [template, layout]
    for value in candidates:
        if value in {'takeaway', 'takeaway + support'}:
            return 'takeaway'
        if value in {'agenda'}:
            return 'agenda'
        if value in {'horizontal process', 'process'}:
            return 'process'
        if value in {'matrix'}:
            return 'matrix'
        if value in {'bar-column', 'bar column'}:
            return 'bar-column'
        if value in {'section-divider', 'section divider'}:
            return 'section-divider'
        if value in {'quote'}:
            return 'quote'
    raise SystemExit(f"unsupported ppt-shapes layout/template for slide id={slide.get('id')}: layout={slide.get('layout')!r} template={slide.get('template')!r}")


def ensure_list(value) -> list:
    if value is None:
        return []
    if isinstance(value, list):
        return value
    return [value]


def compile_slide_payload(spec: dict, slide: dict) -> dict:
    archetype = infer_archetype(slide)
    content = slide.get('content') or {}
    if not isinstance(content, dict):
        raise SystemExit(f"slide content must be a mapping: slide id={slide.get('id')}")

    payload = {
        'archetype': archetype,
        'title': slide.get('title') or storyline_takeaway(spec, slide.get('id')) or 'Untitled slide',
        'subtitle': slide.get('subtitle'),
        'closing': slide.get('closing') or storyline_takeaway(spec, slide.get('id')),
    }

    if archetype == 'takeaway':
        bullets = []
        for lead in ensure_list(content.get('lead')):
            if isinstance(lead, dict):
                bullets.append({'lead': str(lead.get('lead') or lead.get('title') or ''), 'body': str(lead.get('body') or '')})
            else:
                bullets.append({'lead': str(lead), 'body': ''})
        for item in ensure_list(content.get('support')):
            if isinstance(item, dict):
                bullets.append({'lead': str(item.get('lead') or item.get('title') or ''), 'body': str(item.get('body') or '')})
            else:
                bullets.append({'lead': str(item), 'body': ''})
        if not bullets and content.get('body'):
            for item in ensure_list(content.get('body')):
                bullets.append({'lead': str(item), 'body': ''})
        payload['bullets'] = bullets
    elif archetype == 'agenda':
        payload['sections'] = ensure_list(content.get('sections') or content.get('items'))
    elif archetype == 'process':
        payload['phases'] = ensure_list(content.get('steps') or content.get('phases'))
    elif archetype == 'matrix':
        axes = content.get('axes') or {}
        payload['x_axis'] = content.get('x_axis') or axes.get('x')
        payload['y_axis'] = content.get('y_axis') or axes.get('y')
        payload['quadrants'] = ensure_list(content.get('quadrants'))
    elif archetype == 'bar-column':
        payload['bars'] = ensure_list(content.get('bars'))
    elif archetype == 'section-divider':
        payload['section_label'] = content.get('section_label') or slide.get('section_label')
    elif archetype == 'quote':
        quote_value = content.get('quote')
        if isinstance(quote_value, list):
            quote_value = ' '.join(str(x) for x in quote_value)
        payload['quote'] = quote_value or slide.get('quote') or storyline_takeaway(spec, slide.get('id')) or payload['title']
        payload['attribution'] = content.get('attribution') or slide.get('attribution')

    return payload


def clear_existing_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def render_payload(prs: Presentation, payload: dict) -> None:
    archetype = payload['archetype']
    if archetype == 'takeaway':
        add_takeaway_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), bullets=payload.get('bullets', []), closing=payload.get('closing'))
    elif archetype == 'agenda':
        add_agenda_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), sections=payload.get('sections', []), closing=payload.get('closing'))
    elif archetype == 'process':
        add_process_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), phases=payload.get('phases', []), closing=payload.get('closing'))
    elif archetype == 'matrix':
        add_matrix_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), x_axis=payload.get('x_axis'), y_axis=payload.get('y_axis'), quadrants=payload.get('quadrants', []), closing=payload.get('closing'))
    elif archetype == 'bar-column':
        add_bar_column_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), bars=payload.get('bars', []), closing=payload.get('closing'))
    elif archetype == 'section-divider':
        add_section_divider_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), section_label=payload.get('section_label'), closing=payload.get('closing'))
    elif archetype == 'quote':
        add_quote_slide(prs, title=payload['title'], quote=payload['quote'], subtitle=payload.get('subtitle'), attribution=payload.get('attribution'), closing=payload.get('closing'))
    else:
        raise SystemExit(f'unsupported archetype: {archetype}')


def build_from_spec(prs: Presentation, spec: dict) -> None:
    deck_mode = spec.get('deliverable_mode') or spec.get('Deliverable mode')
    slides = spec.get('Slides') or spec.get('slides') or []
    if not isinstance(slides, list):
        raise SystemExit('Slides must be a list in the unified spec sheet')

    rendered = 0
    skipped = 0
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        mode = normalize_mode(deck_mode, slide)
        if mode != 'ppt-shapes':
            skipped += 1
            continue
        payload = compile_slide_payload(spec, slide)
        render_payload(prs, payload)
        rendered += 1

    if rendered == 0:
        raise SystemExit('No ppt-shapes slides found to render from the unified spec sheet')
    if skipped:
        print(f'skipped {skipped} non-ppt-shapes slide(s)', file=sys.stderr)


def default_template_for_spec() -> Path | None:
    return DEFAULT_TEMPLATE if DEFAULT_TEMPLATE.exists() else None


def main() -> int:
    if len(sys.argv) < 2:
        print(
            'usage: build_pptx.py OUTPUT_PPTX [SPEC_YAML_OR_JSON] [TEMPLATE_PPTX]',
            file=sys.stderr,
        )
        return 2

    output = Path(sys.argv[1]).expanduser()
    spec_path = Path(sys.argv[2]).expanduser() if len(sys.argv) > 2 and sys.argv[2] else None
    template_path = Path(sys.argv[3]).expanduser() if len(sys.argv) > 3 and sys.argv[3] else None

    if template_path and template_path.exists():
        template = template_path
    elif spec_path:
        template = default_template_for_spec()
    else:
        template = default_template_for_spec()

    prs = Presentation(str(template)) if template and template.exists() else Presentation()
    if len(prs.slides) > 0:
        clear_existing_slides(prs)

    if spec_path and spec_path.exists():
        structured = load_structured_file(spec_path)
        if 'Slides' in structured or 'slides' in structured:
            build_from_spec(prs, structured)
        elif 'archetype' in structured:
            render_payload(prs, structured)
        else:
            raise SystemExit(f'unrecognized spec structure: {spec_path}')
    elif len(prs.slides) == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = 'Title'
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = 'Subtitle'

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(output)
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
