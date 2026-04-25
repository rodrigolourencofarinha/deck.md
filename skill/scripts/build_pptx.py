#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
import json
import re
import sys

import yaml


BLUE = RGBColor(21, 56, 101)
TEXT = RGBColor(33, 33, 33)
MUTED = RGBColor(68, 68, 68)
PANEL = RGBColor(235, 242, 250)
PANEL_BORDER = RGBColor(210, 222, 238)
WHITE = RGBColor(255, 255, 255)
SKILL_DIR = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = SKILL_DIR / "assets" / "RLF_PPT_Template_v1.pptx"
EMU_PER_INCH = 914400
FRONTMATTER_RE = re.compile(r"\A---\s*\n(.*?)\n---\s*\n", re.DOTALL)
SLIDE_HEADING_RE = re.compile(
    r'^### Slide\s+(.+?)\s+(?:\u2014|-)\s+"(.+?)"\s*$',
    re.MULTILINE,
)
FENCED_YAML_RE = re.compile(r"```yaml\s*\n(.*?)\n```", re.DOTALL)
HTML_COMMENT_RE = re.compile(r"<!--.*?-->", re.DOTALL)


def bool_setting(value, default: bool) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() not in {'0', 'false', 'no', 'off'}
    return bool(value)


def footer_config_from_spec(spec: dict) -> dict:
    production = spec.get('production_defaults') or {}
    if not isinstance(production, dict):
        production = {}
    footer = spec.get('footer') or production.get('footer') or spec.get('Footer') or {}
    if footer is None:
        footer = {}
    if not isinstance(footer, dict):
        footer = {'cr_text': str(footer)}

    return {
        'enabled': bool_setting(footer.get('enabled'), True),
        'page_numbers': bool_setting(footer.get('page_numbers'), True),
        'page_number_format': str(footer.get('page_number_format') or '{page}'),
        'cr_mark': bool_setting(footer.get('cr_mark'), True),
        'cr_text': str(footer.get('cr_text') or footer.get('cr') or 'CR').strip(),
    }


def format_page_number(template: str, page_number: int | str) -> str:
    # Deck standard is standalone slide-order numbering only: 1, 2, 3.
    # Ignore total-count formats such as 1/3 or "1 of 3".
    return str(page_number)


def add_footer_text(slide, x: float, y: float, w: float, h: float, text: str, *, align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = 'Aptos'
    r.font.size = Pt(8)
    r.font.color.rgb = MUTED


def add_standard_footer(
    prs: Presentation,
    slide,
    *,
    page_number: int | str | None,
    footer_config: dict,
) -> None:
    if not footer_config.get('enabled', True):
        return

    slide_w = prs.slide_width / EMU_PER_INCH
    slide_h = prs.slide_height / EMU_PER_INCH
    footer_y = max(slide_h - 0.34, 0.0)

    if footer_config.get('cr_mark', True) and footer_config.get('cr_text'):
        add_footer_text(slide, 0.72, footer_y, 1.4, 0.18, footer_config['cr_text'])

    if footer_config.get('page_numbers', True) and page_number is not None:
        number = format_page_number(footer_config.get('page_number_format') or '{page}', page_number)
        add_footer_text(slide, slide_w - 1.12, footer_y, 0.4, 0.18, number, align=PP_ALIGN.RIGHT)


def add_takeaway_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    bullets: list[dict[str, str]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    panel = slide.shapes.add_shape(1, Inches(0.45), Inches(0.35), Inches(12.35), Inches(6.45))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    top_y = 1.18
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.12), Inches(11.2), Inches(0.68))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(15)
        r.font.italic = True
        r.font.color.rgb = MUTED
        top_y = 1.82

    bullets = bullets or []
    if bullets:
        bullets_box = slide.shapes.add_textbox(Inches(0.9), Inches(top_y), Inches(11.0), Inches(4.05 if closing else 4.7))
        tf = bullets_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        first = True
        for item in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            p.bullet = True
            p.space_after = Pt(6)
            p.line_spacing = 1.05

            head = item.get('lead', '').strip()
            body = item.get('body', '').strip()

            if head:
                r1 = p.add_run()
                r1.text = head
                r1.font.name = 'Aptos'
                r1.font.size = Pt(15)
                r1.font.bold = True
                r1.font.color.rgb = TEXT
            if body:
                r2 = p.add_run()
                prefix = ' ' if head and not body.startswith(' ') else ''
                r2.text = f'{prefix}{body}'
                r2.font.name = 'Aptos'
                r2.font.size = Pt(14)
                r2.font.color.rgb = TEXT

    if closing:
        shape = slide.shapes.add_shape(1, Inches(0.8), Inches(6.08), Inches(11.1), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shape.line.color.rgb = PANEL_BORDER
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.bold = False
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_agenda_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    sections: list[dict[str, str]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    panel = slide.shapes.add_shape(1, Inches(0.45), Inches(0.35), Inches(12.35), Inches(6.45))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    top_y = 1.18
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.12), Inches(11.2), Inches(0.55))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.color.rgb = MUTED
        top_y = 1.75

    sections = sections or []
    start_y = top_y
    box_h = 0.84
    gap = 0.16
    accent_x = 0.95
    content_x = 1.35

    for idx, item in enumerate(sections[:5], start=1):
        y = start_y + (idx - 1) * (box_h + gap)
        accent = slide.shapes.add_shape(1, Inches(accent_x), Inches(y + 0.02), Inches(0.26), Inches(box_h - 0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = BLUE
        accent.line.color.rgb = BLUE

        card = slide.shapes.add_shape(1, Inches(content_x), Inches(y), Inches(10.35), Inches(box_h))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = PANEL_BORDER

        num_box = slide.shapes.add_textbox(Inches(content_x + 0.18), Inches(y + 0.12), Inches(0.55), Inches(0.36))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f'{idx:02d}'
        r.font.name = 'Aptos'
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = BLUE

        head_box = slide.shapes.add_textbox(Inches(content_x + 0.8), Inches(y + 0.09), Inches(8.9), Inches(0.28))
        tf = head_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item.get('title', '')
        r.font.name = 'Aptos'
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT

        body = item.get('body', '').strip()
        if body:
            body_box = slide.shapes.add_textbox(Inches(content_x + 0.8), Inches(y + 0.4), Inches(8.9), Inches(0.22))
            tf = body_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = body
            r.font.name = 'Aptos'
            r.font.size = Pt(11)
            r.font.color.rgb = MUTED

    if closing:
        shape = slide.shapes.add_shape(1, Inches(0.8), Inches(6.12), Inches(11.1), Inches(0.42))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shape.line.color.rgb = PANEL_BORDER
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_process_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    phases: list[dict[str, str]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    subtitle_y = 1.02
    track_y = 1.9
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.08), Inches(11.2), Inches(0.42))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED
        track_y = 2.02

    phases = phases or []
    phases = phases[:6]
    count = max(len(phases), 1)
    left = 0.95
    right = 12.0
    track_w = right - left
    step_gap = 0.12
    box_w = (track_w - step_gap * (count - 1)) / count

    spine = slide.shapes.add_shape(1, Inches(left), Inches(track_y + 0.72), Inches(track_w), Inches(0.1))
    spine.fill.solid()
    spine.fill.fore_color.rgb = PANEL_BORDER
    spine.line.color.rgb = PANEL_BORDER

    for idx, phase in enumerate(phases, start=1):
        x = left + (idx - 1) * (box_w + step_gap)
        circle = slide.shapes.add_shape(1, Inches(x + box_w / 2 - 0.18), Inches(track_y + 0.56), Inches(0.36), Inches(0.36))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.color.rgb = BLUE

        ctf = circle.text_frame
        ctf.clear()
        p = ctf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(idx)
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE

        card = slide.shapes.add_shape(1, Inches(x), Inches(track_y + 1.05), Inches(box_w), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = PANEL_BORDER

        head_box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(track_y + 1.18), Inches(box_w - 0.24), Inches(0.45))
        tf = head_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = phase.get('title', '')
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = TEXT

        body = phase.get('body', '').strip()
        if body:
            body_box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(track_y + 1.67), Inches(box_w - 0.24), Inches(1.1))
            tf = body_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = body
            r.font.name = 'Aptos'
            r.font.size = Pt(11)
            r.font.color.rgb = MUTED

    if closing:
        footer = slide.shapes.add_shape(1, Inches(0.9), Inches(5.85), Inches(11.2), Inches(0.48))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(248, 250, 252)
        footer.line.color.rgb = PANEL_BORDER
        tf = footer.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_matrix_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    x_axis: str | None = None,
    y_axis: str | None = None,
    quadrants: list[dict[str, str]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.2), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    matrix_top = 1.35
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.02), Inches(11.2), Inches(0.38))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED
        matrix_top = 1.65

    matrix_left = 1.45
    matrix_w = 9.9
    matrix_h = 4.2

    outer = slide.shapes.add_shape(1, Inches(matrix_left), Inches(matrix_top), Inches(matrix_w), Inches(matrix_h))
    outer.fill.background()
    outer.line.color.rgb = PANEL_BORDER

    vline = slide.shapes.add_shape(1, Inches(matrix_left + matrix_w / 2), Inches(matrix_top), Inches(0.05), Inches(matrix_h))
    vline.fill.solid()
    vline.fill.fore_color.rgb = PANEL_BORDER
    vline.line.color.rgb = PANEL_BORDER

    hline = slide.shapes.add_shape(1, Inches(matrix_left), Inches(matrix_top + matrix_h / 2), Inches(matrix_w), Inches(0.05))
    hline.fill.solid()
    hline.fill.fore_color.rgb = PANEL_BORDER
    hline.line.color.rgb = PANEL_BORDER

    colors = [RGBColor(248, 250, 252), RGBColor(240, 245, 252), RGBColor(245, 249, 253), RGBColor(237, 243, 250)]
    positions = [
        (matrix_left, matrix_top),
        (matrix_left + matrix_w / 2 + 0.025, matrix_top),
        (matrix_left, matrix_top + matrix_h / 2 + 0.025),
        (matrix_left + matrix_w / 2 + 0.025, matrix_top + matrix_h / 2 + 0.025),
    ]
    quad_w = matrix_w / 2 - 0.025
    quad_h = matrix_h / 2 - 0.025

    quadrants = (quadrants or [])[:4]
    for idx in range(4):
        x, y = positions[idx]
        box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(quad_w), Inches(quad_h))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[idx]
        box.line.color.rgb = WHITE

        item = quadrants[idx] if idx < len(quadrants) else {}
        head_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(quad_w - 0.3), Inches(0.38))
        tf = head_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item.get('title', '')
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = BLUE

        body = item.get('body', '').strip()
        if body:
            body_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.62), Inches(quad_w - 0.3), Inches(1.25))
            tf = body_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = body
            r.font.name = 'Aptos'
            r.font.size = Pt(11)
            r.font.color.rgb = TEXT

    if x_axis:
        x_box = slide.shapes.add_textbox(Inches(matrix_left + 3.1), Inches(matrix_top + matrix_h + 0.12), Inches(3.9), Inches(0.3))
        tf = x_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = x_axis
        r.font.name = 'Aptos'
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED

    if y_axis:
        y_box = slide.shapes.add_textbox(Inches(0.55), Inches(matrix_top + 1.35), Inches(0.7), Inches(1.7))
        tf = y_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = y_axis
        r.font.name = 'Aptos'
        r.font.size = Pt(11)
        r.font.color.rgb = MUTED

    if closing:
        footer = slide.shapes.add_shape(1, Inches(0.9), Inches(6.1), Inches(11.2), Inches(0.42))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(248, 250, 252)
        footer.line.color.rgb = PANEL_BORDER
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_bar_column_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    bars: list[dict[str, str | int | float]] | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE

    chart_top = 1.4
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.08), Inches(11.2), Inches(0.38))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED
        chart_top = 1.72

    bars = bars or []
    max_val = max([float(b.get('value', 0)) for b in bars], default=1.0) or 1.0
    base_x = 1.35
    label_w = 1.9
    bar_x = base_x + label_w
    max_bar_w = 8.2
    row_h = 0.72

    for idx, bar in enumerate(bars[:6]):
        y = chart_top + idx * row_h
        label_box = slide.shapes.add_textbox(Inches(base_x), Inches(y + 0.14), Inches(label_w - 0.15), Inches(0.26))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(bar.get('label', ''))
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = TEXT

        value = float(bar.get('value', 0))
        width = max(0.65, max_bar_w * value / max_val)
        shape = slide.shapes.add_shape(1, Inches(bar_x), Inches(y + 0.08), Inches(width), Inches(0.34))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.color.rgb = BLUE

        value_box = slide.shapes.add_textbox(Inches(bar_x + width + 0.1), Inches(y + 0.11), Inches(0.8), Inches(0.24))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f'{value:g}'
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = BLUE

    if closing:
        footer = slide.shapes.add_shape(1, Inches(0.9), Inches(6.1), Inches(11.2), Inches(0.42))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(248, 250, 252)
        footer.line.color.rgb = PANEL_BORDER
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = BLUE


def add_section_divider_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str | None = None,
    section_label: str | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    bg = slide.shapes.add_shape(1, Inches(0.45), Inches(0.35), Inches(12.35), Inches(6.45))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PANEL
    bg.line.color.rgb = PANEL

    if section_label:
        label_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.2), Inches(2.2), Inches(0.32))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = section_label.upper()
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = BLUE

    title_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.72), Inches(10.6), Inches(1.15))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = BLUE

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.98), Inches(3.08), Inches(10.3), Inches(0.45))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT

    if closing:
        box = slide.shapes.add_textbox(Inches(0.98), Inches(5.55), Inches(10.5), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.italic = True
        r.font.color.rgb = MUTED


def add_quote_slide(
    prs: Presentation,
    *,
    title: str,
    quote: str,
    subtitle: str | None = None,
    attribution: str | None = None,
    closing: str | None = None,
) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.2), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Aptos Display'
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = BLUE

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.2), Inches(0.45))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = 'Aptos'
        r.font.size = Pt(13)
        r.font.color.rgb = MUTED

    quote_card = slide.shapes.add_shape(1, Inches(1.1), Inches(1.75), Inches(10.9), Inches(3.05))
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = PANEL
    quote_card.line.color.rgb = PANEL_BORDER

    q_box = slide.shapes.add_textbox(Inches(1.45), Inches(2.15), Inches(10.2), Inches(2.0))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f'“{quote}”'
    r.font.name = 'Aptos'
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = BLUE

    if attribution:
        a_box = slide.shapes.add_textbox(Inches(7.8), Inches(4.3), Inches(3.4), Inches(0.35))
        tf = a_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = attribution
        r.font.name = 'Aptos'
        r.font.size = Pt(12)
        r.font.italic = True
        r.font.color.rgb = MUTED

    if closing:
        box = slide.shapes.add_textbox(Inches(1.1), Inches(5.3), Inches(10.9), Inches(0.55))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing
        r.font.name = 'Aptos'
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT


def load_yaml_mapping(text: str, path: Path, context: str) -> dict:
    try:
        data = yaml.safe_load(text) or {}
    except yaml.YAMLError as exc:
        raise SystemExit(f'invalid YAML in {context}: {path}: {exc}') from exc
    if not isinstance(data, dict):
        raise SystemExit(f'{context} must be a mapping: {path}')
    return data


def split_markdown_frontmatter(text: str, path: Path) -> tuple[dict, str] | None:
    match = FRONTMATTER_RE.match(text)
    if not match:
        return None
    frontmatter = load_yaml_mapping(match.group(1), path, 'frontmatter')
    return frontmatter, text[match.end():]


def parse_markdown_items(markdown: str) -> list[dict[str, str]]:
    body = FENCED_YAML_RE.sub('', markdown)
    body = HTML_COMMENT_RE.sub('', body)

    items: list[dict[str, str]] = []
    paragraphs: list[str] = []
    in_speaker_notes = False

    for raw_line in body.splitlines():
        line = raw_line.strip()
        if not line or line == '---':
            continue
        lower = line.lower()
        if lower.startswith('**sources:**') or lower.startswith('sources:'):
            in_speaker_notes = False
            continue
        if lower.startswith('**speaker notes:**') or lower.startswith('speaker notes:'):
            in_speaker_notes = True
            continue
        if in_speaker_notes or line.startswith('#') or line.startswith('<'):
            continue

        bullet = re.match(r'^(?:[-*]|\d+[.)])\s+(.+)$', line)
        text = bullet.group(1).strip() if bullet else line
        text = re.sub(r'\s+', ' ', text)

        labeled = re.match(r'^\*\*(.+?)[.:]?\*\*\s*[-:]\s*(.+)$', text)
        if labeled:
            items.append({'lead': labeled.group(1).strip(), 'body': labeled.group(2).strip()})
        elif bullet:
            items.append({'lead': text, 'body': ''})
        else:
            paragraphs.append(text)

    if not items and paragraphs:
        joined = ' '.join(paragraphs).strip()
        if joined:
            items.append({'lead': joined, 'body': ''})
    return items


def parse_deck_markdown(path: Path, text: str) -> dict | None:
    split = split_markdown_frontmatter(text, path)
    if not split:
        return None

    frontmatter, body = split
    if '## Slides' not in body or 'deck' not in frontmatter:
        return None

    slides_section = body.split('## Slides', 1)[1]
    headings = list(SLIDE_HEADING_RE.finditer(slides_section))
    if not headings:
        raise SystemExit(f'No deck.md slide headings found in {path}')

    slides: list[dict] = []
    for index, heading in enumerate(headings):
        segment_end = headings[index + 1].start() if index + 1 < len(headings) else len(slides_section)
        segment = slides_section[heading.end():segment_end]

        metadata: dict = {}
        for block_index, block in enumerate(FENCED_YAML_RE.findall(segment)):
            block_data = load_yaml_mapping(block, path, f'slide {index + 1} YAML block')
            if block_index == 0 and {'id', 'type', 'layout', 'mode', 'image_decision'} & set(block_data):
                metadata.update(block_data)
            else:
                for key in ('chart', 'creative_direction', 'required_text'):
                    if key in block_data:
                        metadata[key] = block_data[key]

        slide_type = str(metadata.get('type') or '').strip()
        body_items = parse_markdown_items(segment)
        content: dict = {'body': body_items}
        if slide_type == 'roadmap':
            content['phases'] = body_items
        if slide_type == 'framework' and body_items:
            content['quadrants'] = body_items[:4]
        if metadata.get('chart'):
            content['chart'] = metadata['chart']

        slide = {
            'id': metadata.get('id') or heading.group(1).strip(),
            'title': heading.group(2).strip(),
            'type': slide_type,
            'layout': metadata.get('layout'),
            'mode': metadata.get('mode'),
            'slide_mode': metadata.get('mode'),
            'template': metadata.get('template'),
            'image_decision': metadata.get('image_decision'),
            'content': content,
        }
        if metadata.get('chart'):
            slide['chart'] = metadata['chart']
        slides.append(slide)

    return {
        'schema_version': frontmatter.get('schema_version'),
        'deck': frontmatter.get('deck') or {},
        'narrative_template': frontmatter.get('narrative_template'),
        'production_defaults': frontmatter.get('production_defaults') or {},
        'image_generation': frontmatter.get('image_generation') or {},
        'design_tokens': frontmatter.get('design_tokens') or {},
        'slides': slides,
    }


def load_structured_file(path: Path):
    text = path.read_text(encoding='utf-8')
    if path.suffix.lower() == '.md':
        deck = parse_deck_markdown(path, text)
        if deck is not None:
            return deck
        return load_yaml_mapping(text, path, 'spec root')
    if path.suffix.lower() in {'.yaml', '.yml'}:
        return load_yaml_mapping(text, path, 'spec root')
    return json.loads(text)


def slugify(value: str) -> str:
    value = value.lower().strip()
    value = re.sub(r'[^a-z0-9]+', '-', value)
    return value.strip('-') or 'deck'


def normalize_mode(deck_mode: str | None, slide: dict) -> str:
    deck_mode = (deck_mode or '').strip().lower()
    slide_mode = str(slide.get('slide_mode') or slide.get('mode') or '').strip().lower()
    if deck_mode == 'mixed':
        if slide_mode not in {'ppt-shapes', 'designer-mode'}:
            raise SystemExit(f"mixed decks require slide_mode per slide; missing on slide id={slide.get('id')}")
        return slide_mode
    if slide_mode in {'ppt-shapes', 'designer-mode'}:
        return slide_mode
    if deck_mode in {'ppt-shapes', 'designer-mode'}:
        return deck_mode
    return slide_mode or 'ppt-shapes'


def storyline_takeaway(spec: dict, slide_id: int | str | None) -> str | None:
    for item in spec.get('Storyline', []) or []:
        if not isinstance(item, dict):
            continue
        if str(item.get('slide')) == str(slide_id):
            return item.get('key_takeaway')
    return None


def infer_archetype(slide: dict) -> str:
    slide_type = str(slide.get('type') or '').strip().lower().replace('_', '-')
    template = str(slide.get('template') or '').strip().lower()
    layout = str(slide.get('layout') or '').strip().lower()
    content = slide.get('content') or {}
    chart = slide.get('chart') or (content.get('chart') if isinstance(content, dict) else None) or {}
    chart_type = str(chart.get('type') or '').strip().lower() if isinstance(chart, dict) else ''
    candidates = [template, layout, slide_type]
    for value in candidates:
        if value in {'takeaway', 'takeaway + support'}:
            return 'takeaway'
        if value in {'agenda'}:
            return 'agenda'
        if value in {'horizontal process', 'process'}:
            return 'process'
        if value in {'matrix'}:
            return 'matrix'
        if value in {'bar-column', 'bar column'}:
            return 'bar-column'
        if value in {'section-divider', 'section divider'}:
            return 'section-divider'
        if value in {'quote'}:
            return 'quote'
    if slide_type == 'section-divider':
        return 'section-divider'
    if slide_type == 'roadmap' or 'timeline' in layout or 'roadmap' in layout:
        return 'process'
    if slide_type == 'framework' and any(token in layout for token in ('2x2', 'matrix', 'quadrant')):
        return 'matrix'
    if slide_type == 'chart' or chart_type:
        if any(token in chart_type for token in ('bar', 'column')):
            return 'bar-column'
        return 'takeaway'
    if slide_type in {
        'executive-summary',
        'situation',
        'complication',
        'key-takeaways',
        'analysis',
        'framework',
        'recommendation',
        'risk-mitigation',
        'next-steps',
        'appendix',
    }:
        return 'takeaway'
    raise SystemExit(f"unsupported ppt-shapes layout/template/type for slide id={slide.get('id')}: layout={slide.get('layout')!r} template={slide.get('template')!r} type={slide.get('type')!r}")


def ensure_list(value) -> list:
    if value is None:
        return []
    if isinstance(value, list):
        return value
    return [value]


def append_bullet(bullets: list[dict[str, str]], item) -> None:
    if isinstance(item, dict):
        bullets.append({
            'lead': str(item.get('lead') or item.get('title') or item.get('label') or ''),
            'body': str(item.get('body') or item.get('caption') or item.get('description') or ''),
        })
    else:
        bullets.append({'lead': str(item), 'body': ''})


def compile_slide_payload(spec: dict, slide: dict) -> dict:
    archetype = infer_archetype(slide)
    content = slide.get('content') or {}
    if not isinstance(content, dict):
        raise SystemExit(f"slide content must be a mapping: slide id={slide.get('id')}")

    payload = {
        'archetype': archetype,
        'title': slide.get('title') or storyline_takeaway(spec, slide.get('id')) or 'Untitled slide',
        'subtitle': slide.get('subtitle'),
        'closing': slide.get('closing') or storyline_takeaway(spec, slide.get('id')),
    }

    if archetype == 'takeaway':
        bullets = []
        for lead in ensure_list(content.get('lead')):
            append_bullet(bullets, lead)
        for item in ensure_list(content.get('support')):
            append_bullet(bullets, item)
        if not bullets and content.get('body'):
            for item in ensure_list(content.get('body')):
                append_bullet(bullets, item)
        payload['bullets'] = bullets
    elif archetype == 'agenda':
        payload['sections'] = ensure_list(content.get('sections') or content.get('items'))
    elif archetype == 'process':
        payload['phases'] = ensure_list(content.get('steps') or content.get('phases') or content.get('body'))
    elif archetype == 'matrix':
        axes = content.get('axes') or {}
        payload['x_axis'] = content.get('x_axis') or axes.get('x')
        payload['y_axis'] = content.get('y_axis') or axes.get('y')
        payload['quadrants'] = ensure_list(content.get('quadrants') or content.get('body'))
    elif archetype == 'bar-column':
        payload['bars'] = ensure_list(content.get('bars'))
    elif archetype == 'section-divider':
        payload['section_label'] = content.get('section_label') or slide.get('section_label')
    elif archetype == 'quote':
        quote_value = content.get('quote')
        if isinstance(quote_value, list):
            quote_value = ' '.join(str(x) for x in quote_value)
        payload['quote'] = quote_value or slide.get('quote') or storyline_takeaway(spec, slide.get('id')) or payload['title']
        payload['attribution'] = content.get('attribution') or slide.get('attribution')

    return payload


def clear_existing_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def render_payload(
    prs: Presentation,
    payload: dict,
    *,
    page_number: int | str | None = None,
    footer_config: dict | None = None,
) -> None:
    archetype = payload['archetype']
    before_count = len(prs.slides)
    if archetype == 'takeaway':
        add_takeaway_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), bullets=payload.get('bullets', []), closing=payload.get('closing'))
    elif archetype == 'agenda':
        add_agenda_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), sections=payload.get('sections', []), closing=payload.get('closing'))
    elif archetype == 'process':
        add_process_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), phases=payload.get('phases', []), closing=payload.get('closing'))
    elif archetype == 'matrix':
        add_matrix_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), x_axis=payload.get('x_axis'), y_axis=payload.get('y_axis'), quadrants=payload.get('quadrants', []), closing=payload.get('closing'))
    elif archetype == 'bar-column':
        add_bar_column_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), bars=payload.get('bars', []), closing=payload.get('closing'))
    elif archetype == 'section-divider':
        add_section_divider_slide(prs, title=payload['title'], subtitle=payload.get('subtitle'), section_label=payload.get('section_label'), closing=payload.get('closing'))
    elif archetype == 'quote':
        add_quote_slide(prs, title=payload['title'], quote=payload['quote'], subtitle=payload.get('subtitle'), attribution=payload.get('attribution'), closing=payload.get('closing'))
    else:
        raise SystemExit(f'unsupported archetype: {archetype}')

    if footer_config is not None and len(prs.slides) > before_count:
        add_standard_footer(
            prs,
            prs.slides[len(prs.slides) - 1],
            page_number=page_number,
            footer_config=footer_config,
        )


def build_from_spec(prs: Presentation, spec: dict) -> None:
    production_defaults = spec.get('production_defaults') or {}
    deck_mode = (
        production_defaults.get('default_slide_mode')
        or spec.get('deliverable_mode')
        or spec.get('Deliverable mode')
    )
    slides = spec.get('Slides') or spec.get('slides') or []
    if not isinstance(slides, list):
        raise SystemExit('Slides must be a list in the deck spec')
    footer_config = footer_config_from_spec(spec)

    rendered = 0
    skipped = 0
    for page_number, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        mode = normalize_mode(deck_mode, slide)
        if mode != 'ppt-shapes':
            skipped += 1
            continue
        payload = compile_slide_payload(spec, slide)
        render_payload(prs, payload, page_number=page_number, footer_config=footer_config)
        rendered += 1

    if rendered == 0:
        raise SystemExit('No ppt-shapes slides found to render from the deck spec')
    if skipped:
        print(f'skipped {skipped} non-ppt-shapes slide(s)', file=sys.stderr)


def default_template_for_spec() -> Path | None:
    return DEFAULT_TEMPLATE if DEFAULT_TEMPLATE.exists() else None


def main() -> int:
    if len(sys.argv) < 2:
        print(
            'usage: build_pptx.py OUTPUT_PPTX [DECK_MD_OR_SPEC_YAML_OR_JSON] [TEMPLATE_PPTX]',
            file=sys.stderr,
        )
        return 2

    output = Path(sys.argv[1]).expanduser()
    spec_path = Path(sys.argv[2]).expanduser() if len(sys.argv) > 2 and sys.argv[2] else None
    template_path = Path(sys.argv[3]).expanduser() if len(sys.argv) > 3 and sys.argv[3] else None

    if template_path and template_path.exists():
        template = template_path
    elif spec_path:
        template = default_template_for_spec()
    else:
        template = default_template_for_spec()

    prs = Presentation(str(template)) if template and template.exists() else Presentation()
    if len(prs.slides) > 0:
        clear_existing_slides(prs)

    if spec_path and spec_path.exists():
        structured = load_structured_file(spec_path)
        if 'Slides' in structured or 'slides' in structured:
            build_from_spec(prs, structured)
        elif 'archetype' in structured:
            render_payload(prs, structured, page_number=1, footer_config=footer_config_from_spec(structured))
        else:
            raise SystemExit(f'unrecognized spec structure: {spec_path}')
    elif len(prs.slides) == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = 'Title'
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = 'Subtitle'

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(output)
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
