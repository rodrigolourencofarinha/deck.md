#!/usr/bin/env python3
from __future__ import annotations

import argparse
from collections import Counter
from pathlib import Path
import posixpath
import re
import unicodedata
from zipfile import ZIP_DEFLATED, ZipFile
import xml.etree.ElementTree as ET

from pptx import Presentation


SKILL_DIR = Path(__file__).resolve().parent.parent
DEFAULT_SOURCE = SKILL_DIR / "assets" / "templates_full.pptx"
DEFAULT_OUTPUT_DIR = SKILL_DIR / "assets" / "templates"
DEFAULT_CATALOG = SKILL_DIR / "references" / "template-catalog.md"

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
SLIDE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"

ET.register_namespace("p", P_NS)
ET.register_namespace("r", R_NS)
ET.register_namespace("", PKG_REL_NS)


def clean_text(text: str) -> str:
    return " ".join(text.split()).strip()


def slide_title(slide, index: int) -> str:
    if slide.shapes.title is not None:
        title = clean_text(slide.shapes.title.text)
        if title:
            return title

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = clean_text(shape.text)
            if text and not text.isdigit():
                return text[:80]

    return f"Slide {index:03d}"


def slugify(text: str) -> str:
    text = text.replace("\u00d7", "x").replace("&", " and ")
    normalized = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    normalized = normalized.lower()
    normalized = re.sub(r"[^a-z0-9]+", "-", normalized)
    normalized = re.sub(r"-+", "-", normalized).strip("-")
    return normalized or "untitled"


def classify(title: str) -> tuple[str, str]:
    t = title.lower()

    if "title page" in t:
        return ("Cover", "Opening cover with title, subtitle, author, and date.")
    if t == "only title":
        return ("Divider", "Simple section divider or transition slide.")
    if "dashboard" in t:
        return ("Dashboard", "Multiple metrics, panels, or status blocks on one slide.")
    if "agenda" in t or "table of contents" in t or "schedule" in t:
        return ("Agenda / schedule", "Agenda, table of contents, meeting flow, or timed session plan.")
    if "picture" in t or "image" in t:
        return ("Image cards", "Text-and-image card layout for examples, cases, people, or offerings.")
    if "column" in t or "text box" in t:
        return ("Text blocks", "Parallel text blocks for comparing ideas, options, or workstreams.")
    if "goal" in t or "objective" in t or "assessment" in t or "building block" in t:
        return ("Goals / building blocks", "Objectives, capability blocks, assessment criteria, or priorities.")
    if "matrix" in t or "matrices" in t or "swot" in t or "prioritization" in t:
        return ("Matrix / prioritization", "Two-axis, multi-cell, SWOT, portfolio, or prioritization analysis.")
    if "funnel" in t:
        return ("Funnel", "Conversion funnel, staged narrowing, pipeline, or selection logic.")
    if "pyramid" in t or "ziggurat" in t or "layer" in t or "triangle" in t:
        return ("Hierarchy / layers", "Layered hierarchy, pyramid, maturity, foundation, or stack concept.")
    if "cause and effect" in t or "influencing" in t or "component" in t:
        return ("Cause / system", "Drivers, effects, components, dependencies, or system logic.")
    if "journey" in t:
        return ("Journey", "Customer journey, touchpoints, or experience map.")
    if "swim lane" in t:
        return ("Swim lane", "Cross-functional process by owner, lane, team, or phase.")
    if "value chain" in t:
        return ("Value chain", "Sequential business activities, operating model, or value-chain logic.")
    if "core elements" in t:
        return ("Core elements", "Three-part framework or central strategic pillars.")
    if (
        "process" in t
        or "phase" in t
        or "stage" in t
        or "steps" in t
        or "path" in t
        or "sprint" in t
        or "circle" in t
        or "circular" in t
        or "cycle" in t
        or "flow" in t
        or "plan" in t
    ):
        return ("Process / flow", "Steps, phases, roadmap, loop, sequence, or workflow structure.")

    return ("General framework", "Reusable consulting-style structure for adapting to a slide message.")


def escape_md(text: str) -> str:
    return text.replace("|", "\\|")


def xml_bytes(root: ET.Element) -> bytes:
    return b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n' + ET.tostring(root, encoding="utf-8")


def rels_path_for(part_name: str) -> str:
    directory, filename = posixpath.split(part_name)
    if directory:
        return posixpath.join(directory, "_rels", filename + ".rels")
    return posixpath.join("_rels", filename + ".rels")


def write_one_slide_package(
    zin: ZipFile,
    names: set[str],
    rid: str,
    dest: Path,
) -> None:
    one_pres_root = ET.fromstring(zin.read("ppt/presentation.xml"))
    one_sld_id_lst = one_pres_root.find(f"{{{P_NS}}}sldIdLst")
    if one_sld_id_lst is None:
        raise RuntimeError("No slide ID list found in ppt/presentation.xml")

    for child in list(one_sld_id_lst):
        if child.get(f"{{{R_NS}}}id") != rid:
            one_sld_id_lst.remove(child)

    one_pres_rels_root = ET.fromstring(zin.read("ppt/_rels/presentation.xml.rels"))
    selected_slide_part = ""
    for rel in list(one_pres_rels_root):
        if rel.get("Type") == SLIDE_REL:
            if rel.get("Id") == rid:
                selected_slide_part = posixpath.normpath("ppt/" + str(rel.get("Target", "")))
            else:
                one_pres_rels_root.remove(rel)

    if not selected_slide_part:
        raise RuntimeError(f"Could not find selected slide relationship {rid}")

    selected_slide_rels = rels_path_for(selected_slide_part)
    modified = {
        "ppt/presentation.xml": xml_bytes(one_pres_root),
        "ppt/_rels/presentation.xml.rels": xml_bytes(one_pres_rels_root),
    }

    with ZipFile(dest, "w", ZIP_DEFLATED) as zout:
        for name in sorted(names):
            if re.match(r"ppt/slides/slide\d+\.xml$", name) and name != selected_slide_part:
                continue
            if re.match(r"ppt/slides/_rels/slide\d+\.xml\.rels$", name) and name != selected_slide_rels:
                continue
            data = modified.get(name) or zin.read(name)
            zout.writestr(name, data)


def build_catalog(rows: list[dict[str, str | int]], catalog: Path) -> None:
    lines = [
        "# Template Slide Catalog",
        "",
        "This catalog indexes the one-slide PowerPoint templates extracted from `assets/templates_full.pptx`.",
        "",
        "Use these files as structural references when building editable consulting-style slides. The files live in `assets/templates/` and each PPTX contains one visible slide while preserving the relevant PowerPoint theme, layouts, media, and editable objects.",
        "",
        "Selection rule: choose the template by slide job first, then adapt text and visual hierarchy to the current message. Do not copy placeholder lorem ipsum text into final work.",
        "",
        "| # | Template | Title | Category | Best use |",
        "|---:|---|---|---|---|",
    ]

    for row in rows:
        rel_link = f"../assets/templates/{row['file']}"
        lines.append(
            f"| {row['slide']} | [`{row['file']}`]({rel_link}) | {escape_md(str(row['title']))} | {escape_md(str(row['category']))} | {escape_md(str(row['use']))} |"
        )

    catalog.write_text("\n".join(lines) + "\n", encoding="utf-8")


def split_templates(source: Path, output_dir: Path, catalog: Path) -> int:
    output_dir.mkdir(parents=True, exist_ok=True)
    for old in output_dir.glob("*.pptx"):
        old.unlink()

    prs = Presentation(str(source))
    titles = [slide_title(slide, idx) for idx, slide in enumerate(prs.slides, 1)]

    rows: list[dict[str, str | int]] = []
    slug_counts: Counter[str] = Counter()

    with ZipFile(source) as zin:
        names = set(zin.namelist())
        pres_root = ET.fromstring(zin.read("ppt/presentation.xml"))
        sld_id_lst = pres_root.find(f"{{{P_NS}}}sldIdLst")
        if sld_id_lst is None:
            raise RuntimeError("No slide ID list found in ppt/presentation.xml")
        slide_ids = list(sld_id_lst)

        if len(slide_ids) != len(titles):
            raise RuntimeError(f"Slide count mismatch: package has {len(slide_ids)}, python-pptx has {len(titles)}")

        for idx, slide_id in enumerate(slide_ids, 1):
            rid = slide_id.get(f"{{{R_NS}}}id")
            if not rid:
                raise RuntimeError(f"Slide {idx} has no relationship id")

            title = titles[idx - 1]
            slug = slugify(title)
            slug_counts[slug] += 1
            if slug_counts[slug] > 1:
                slug = f"{slug}-{slug_counts[slug]}"
            filename = f"{idx:03d}-{slug}.pptx"
            write_one_slide_package(zin, names, rid, output_dir / filename)

            category, use = classify(title)
            rows.append({"slide": idx, "file": filename, "title": title, "category": category, "use": use})

    build_catalog(rows, catalog)
    return len(rows)


def main() -> int:
    parser = argparse.ArgumentParser(description="Split templates_full.pptx into one-slide template PPTX files.")
    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE)
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT_DIR)
    parser.add_argument("--catalog", type=Path, default=DEFAULT_CATALOG)
    args = parser.parse_args()

    count = split_templates(args.source, args.output_dir, args.catalog)
    print(f"Created {count} template PPTX files in {args.output_dir}")
    print(f"Wrote catalog: {args.catalog}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
